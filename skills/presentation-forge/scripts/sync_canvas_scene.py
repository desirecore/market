#!/usr/bin/env python3
"""Extract editable PowerPoint objects into unified scene element v2 canvas records."""

from __future__ import annotations

import argparse
import hashlib
import json
import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN


SLIDE_W = 13.333333
SLIDE_H = 7.5
CANVAS_W = 1920
CANVAS_H = 1080


def load(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def save(path: Path, value: object) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(value, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def canonical_hash(value: object) -> str:
    raw = json.dumps(value, ensure_ascii=False, sort_keys=True, separators=(",", ":")).encode()
    return hashlib.sha256(raw).hexdigest()


def object_id(shape: Any) -> str:
    return str(getattr(shape, "name", "")).split(" [pf-role=", 1)[0]


def object_role(shape: Any) -> str | None:
    name = str(getattr(shape, "name", ""))
    return name.split("[pf-role=", 1)[1].rstrip("]") if "[pf-role=" in name else None


def rgb_value(color: Any) -> str | None:
    try:
        rgb = color.rgb
        return str(rgb) if rgb is not None else None
    except (AttributeError, TypeError, ValueError):
        return None


def first_run(shape: Any) -> Any | None:
    if not getattr(shape, "has_text_frame", False):
        return None
    for paragraph in shape.text_frame.paragraphs:
        if paragraph.runs:
            return paragraph.runs[0]
    return None


def shape_style(shape: Any) -> dict[str, Any]:
    style: dict[str, Any] = {}
    try:
        fill = rgb_value(shape.fill.fore_color)
        if fill:
            style["fill"] = fill
    except (AttributeError, TypeError):
        pass
    try:
        line = rgb_value(shape.line.color)
        if line:
            style["line"] = line
        if shape.line.width:
            style["line_width_pt"] = round(float(shape.line.width.pt), 2)
    except (AttributeError, TypeError):
        pass
    run = first_run(shape)
    if run is not None:
        if run.font.size:
            style["font_size_pt"] = round(float(run.font.size.pt), 2)
        if run.font.name:
            style["font_family"] = run.font.name
        font_color = rgb_value(run.font.color)
        if font_color:
            style["font_color"] = font_color
        if run.font.bold is not None:
            style["bold"] = bool(run.font.bold)
        paragraph = run._parent
        style["align"] = {
            PP_ALIGN.CENTER: "center", PP_ALIGN.RIGHT: "right", PP_ALIGN.JUSTIFY: "justify",
        }.get(paragraph.alignment, "left")
        try:
            anchor = shape.text_frame.vertical_anchor
            style["vertical_align"] = {1: "top", 3: "middle", 4: "bottom"}.get(int(anchor), "middle") if anchor is not None else "middle"
        except (AttributeError, TypeError, ValueError):
            pass
    return style


def shape_kind(shape: Any) -> str:
    if getattr(shape, "has_table", False):
        return "table"
    if getattr(shape, "has_chart", False):
        return "chart"
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    if shape.shape_type == MSO_SHAPE_TYPE.LINE:
        return "connector"
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return "group"
    if getattr(shape, "has_text_frame", False) and str(shape.text).strip():
        return "text"
    return "shape"


def shape_geometry(shape: Any) -> str | None:
    try:
        value = shape.auto_shape_type
        name = getattr(value, "name", None)
        if name:
            return str(name).lower()
        text = str(value).split(" ", 1)[0].lower()
        return text if text and text != "none" else None
    except (AttributeError, TypeError, ValueError):
        return None


def source_binding(slide_index: int, shape_id: str) -> dict[str, str] | None:
    suffix = re.sub(r"^s\d{3}_", "", shape_id)
    direct = {
        "hero": "title", "title": "title", "closing_title": "title",
        "subtitle": "subtitle", "message": "key_message", "closing_message": "key_message",
    }
    if suffix in direct:
        return {"document": "analysis/native_deck_spec.json", "path": f"/slides/{slide_index}/{direct[suffix]}"}
    patterns = (
        (r"card_(\d+)_(title|body)$", "cards"),
        (r"step_(\d+)_(title|body)$", "steps"),
        (r"metric_(\d+)_(value|label)$", "metrics"),
        (r"layer_(\d+)_(title|body)$", "layers"),
    )
    for pattern, collection in patterns:
        match = re.fullmatch(pattern, suffix)
        if match:
            return {"document": "analysis/native_deck_spec.json", "path": f"/slides/{slide_index}/{collection}/{int(match.group(1)) - 1}/{match.group(2)}"}
    match = re.fullmatch(r"compare_(\d+)_title", suffix)
    if match:
        side = "left" if int(match.group(1)) == 1 else "right"
        return {"document": "analysis/native_deck_spec.json", "path": f"/slides/{slide_index}/{side}/title"}
    match = re.fullmatch(r"action_(\d+)_text", suffix)
    if match:
        return {"document": "analysis/native_deck_spec.json", "path": f"/slides/{slide_index}/actions/{int(match.group(1)) - 1}"}
    return None


def capabilities(kind: str) -> list[str]:
    if kind == "text":
        return ["text", "geometry", "style", "rotation", "z-order"]
    if kind in {"shape", "connector"}:
        return ["geometry", "style", "rotation", "z-order"]
    if kind in {"image", "chart", "table", "group"}:
        return ["geometry", "rotation", "z-order"]
    return ["geometry"]


def element_from_shape(shape: Any, slide_index: int, z_index: int) -> dict[str, Any] | None:
    shape_id = object_id(shape)
    if not re.fullmatch(r"s\d{3}_.+", shape_id):
        return None
    x = float(shape.left.inches) / SLIDE_W * CANVAS_W
    y = float(shape.top.inches) / SLIDE_H * CANVAS_H
    w = max(1.0, float(shape.width.inches) / SLIDE_W * CANVAS_W)
    h = max(1.0, float(shape.height.inches) / SLIDE_H * CANVAS_H)
    kind = shape_kind(shape)
    element: dict[str, Any] = {
        "id": shape_id,
        "type": kind,
        "bbox": [round(x, 2), round(y, 2), round(w, 2), round(h, 2)],
        "rotation": round(float(getattr(shape, "rotation", 0) or 0), 2),
        "z_index": z_index,
        "editable": True,
        "locked": False,
        "capabilities": capabilities(kind),
        "role": object_role(shape),
        "style": shape_style(shape),
    }
    if kind == "text":
        element["text"] = str(shape.text)
    geometry = shape_geometry(shape)
    if geometry:
        element["geometry"] = geometry
    binding = source_binding(slide_index, shape_id)
    if binding:
        element["source_binding"] = binding
    return element


def refresh_hash(scene: dict[str, Any]) -> None:
    hashable = dict(scene)
    hashable.pop("revision", None)
    dependencies = dict(scene.get("dependencies", {})); dependencies.pop("scene_hash", None)
    hashable["dependencies"] = dependencies
    scene.setdefault("dependencies", {})["scene_hash"] = canonical_hash(hashable)


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--session", required=True)
    parser.add_argument("--pptx", required=True)
    parser.add_argument("--report")
    args = parser.parse_args()
    session = Path(args.session).resolve()
    pptx = Path(args.pptx).resolve()
    presentation = Presentation(pptx)
    scenes_dir = session / "scenes"; scenes_dir.mkdir(parents=True, exist_ok=True)
    counts: list[dict[str, int]] = []
    for slide_index, slide in enumerate(presentation.slides):
        number = slide_index + 1
        scene_path = scenes_dir / f"slide-{number:03d}.scene.json"
        if not scene_path.is_file():
            raise SystemExit(f"scene missing: {scene_path}; run compile_scenes.py first")
        scene = load(scene_path)
        elements = [element for z, shape in enumerate(slide.shapes) if (element := element_from_shape(shape, slide_index, z))]
        previous = canonical_hash(scene.get("elements", []))
        scene["schema_version"] = 2
        scene["canvas"] = {"width": CANVAS_W, "height": CANVAS_H, "unit": "px"}
        scene["elements"] = elements
        if previous != canonical_hash(elements):
            scene["revision"] = int(scene.get("revision", 1)) + 1
        refresh_hash(scene)
        save(scene_path, scene)
        counts.append({"slide_number": number, "element_count": len(elements)})
    report = {"schema_version": 1, "status": "PASS", "pptx": str(pptx), "slide_count": len(counts), "slides": counts}
    report_path = Path(args.report).resolve() if args.report else session / "reports" / "canvas-sync.json"
    save(report_path, report)
    print(json.dumps(report, ensure_ascii=False))


if __name__ == "__main__":
    main()
