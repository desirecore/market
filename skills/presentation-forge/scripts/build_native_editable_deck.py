#!/usr/bin/env python3
"""Build an AIPPT-style native editable deck with optional image-2 visual assets."""

from __future__ import annotations

import argparse
import hashlib
import json
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from build_semantic_deck import fit_contain, hex_to_rgb, patch_theme_east_asian_fonts, set_run_typefaces
from style_presets import merged_design_tokens, resolve_style
from validate_design_quality import analyze_design


ROOT = Path(__file__).resolve().parents[1]
SLIDE_W = 13.333333
SLIDE_H = 7.5
POLICIES = {"native-only", "native-image-assisted", "image-led-editable"}
ARCHETYPES = {"cover", "content-structured", "process-flow", "comparison-two-zone", "data-callouts", "table", "architecture", "closing-action"}
ROLE_SIZE_KEY = {
    "hero": "hero",
    "section_title": "section_title",
    "title": "page_title",
    "subtitle": "subtitle",
    "header": "minor_title",
    "body": "body",
    "label": "label",
    "caption": "caption",
    "table": "table",
}


def read_json(path: Path) -> dict[str, Any]:
    value = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(value, dict):
        raise ValueError(f"expected JSON object: {path}")
    return value


def sha256(path: Path) -> str:
    digest = hashlib.sha256()
    with path.open("rb") as handle:
        for block in iter(lambda: handle.read(1024 * 1024), b""):
            digest.update(block)
    return digest.hexdigest()


def first_color(value: Any, fallback: str) -> str:
    if isinstance(value, list):
        return str(value[0]) if value else fallback
    return str(value or fallback)


def second_color(value: Any, fallback: str) -> str:
    if isinstance(value, list) and len(value) > 1:
        return str(value[1])
    return fallback


def load_design(style_id: str, variant_id: str | None, profile_id: str | None, table_id: str | None) -> tuple[dict[str, str], dict[str, Any], dict[str, Any], str, str, str]:
    style, variant, raw = merged_design_tokens(style_id, variant_id)
    text_values = raw.get("text") if isinstance(raw.get("text"), list) else [raw.get("text")]
    background = first_color(raw.get("background"), "FFFFFF")
    colors = {
        "background": background,
        "primary": first_color(raw.get("primary"), "2563EB"),
        "secondary": first_color(raw.get("secondary"), "60A5FA"),
        "text": str((text_values or ["111827"])[0] or "111827"),
        "muted": str((text_values or ["111827", "64748B"])[1] if len(text_values or []) > 1 else "64748B"),
        "panel": second_color(raw.get("background"), "F3F6FA"),
    }
    profiles = read_json(ROOT / "styles" / "typography-profiles.json")
    resolved_profile = profile_id or style.get("typography_profile")
    resolved_table = table_id or style.get("table_profile")
    typography = next((row for row in profiles.get("typography_profiles", []) if row.get("id") == resolved_profile), None)
    table_policy = next((row for row in profiles.get("table_profiles", []) if row.get("id") == resolved_table), None)
    if not typography or not table_policy:
        raise ValueError(f"unknown typography/table profile: {resolved_profile}/{resolved_table}")
    return colors, typography, table_policy, str(resolved_profile), str(resolved_table), str(variant["id"])


def set_role(shape: Any, role: str, object_id: str) -> None:
    shape._element.nvSpPr.cNvPr.set("name", f"{object_id} [pf-role={role}]")


def set_object_name(shape: Any, object_id: str) -> None:
    for attr in ("nvSpPr", "nvCxnSpPr", "nvPicPr", "nvGraphicFramePr", "nvGrpSpPr"):
        node = getattr(shape._element, attr, None)
        if node is not None and getattr(node, "cNvPr", None) is not None:
            node.cNvPr.set("name", object_id)
            return


def object_id(shape: Any) -> str:
    return str(getattr(shape, "name", "")).split(" [pf-role=", 1)[0]


def object_role(shape: Any) -> str | None:
    name = str(getattr(shape, "name", ""))
    return name.split("[pf-role=", 1)[1].rstrip("]") if "[pf-role=" in name else None


def role_settings(profile: dict[str, Any], role: str) -> tuple[float, str, str, dict[str, Any]]:
    size_key = ROLE_SIZE_KEY[role]
    size = float(profile["tokens"][size_key])
    heading = size_key in {"hero", "section_title", "page_title", "subtitle", "minor_title"}
    fonts = profile["fonts"]
    east_asia = fonts["heading_east_asia"] if heading else fonts["body_east_asia"]
    latin = fonts["latin"]
    paragraph_group = "title" if heading else "caption" if size_key in {"label", "caption"} else "body"
    return size, east_asia, latin, profile.get("paragraph", {}).get(paragraph_group, {})


def add_text(
    slide: Any,
    text: str,
    box: tuple[float, float, float, float],
    role: str,
    profile: dict[str, Any],
    color: str,
    object_id: str,
    *,
    align: str = "left",
    bold: bool | None = None,
    vertical: str = "middle",
) -> Any:
    x, y, w, h = box
    shape = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    set_role(shape, role, object_id)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.04)
    frame.margin_right = Inches(0.04)
    frame.margin_top = Inches(0.02)
    frame.margin_bottom = Inches(0.02)
    shape.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}.get(vertical, MSO_ANCHOR.MIDDLE)
    size, east_asia, latin, paragraph_policy = role_settings(profile, role)
    lines = str(text).split("\n")
    for index, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        run = paragraph.add_run()
        run.text = line
        run.font.name = latin
        run.font.size = Pt(size)
        run.font.bold = bool(role in {"hero", "section_title", "title", "header"}) if bold is None else bold
        run.font.color.rgb = hex_to_rgb(color, "111827")
        set_run_typefaces(run, east_asia, latin, profile["fonts"].get("complex_script") or latin)
        paragraph.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(align, PP_ALIGN.LEFT)
        paragraph.line_spacing = float(paragraph_policy.get("line_spacing_multiple", 1.0))
        paragraph.space_before = Pt(size * float(paragraph_policy.get("space_before_lines", 0)))
        paragraph.space_after = Pt(size * float(paragraph_policy.get("space_after_lines", 0)))
    return shape


def add_box(slide: Any, box: tuple[float, float, float, float], fill: str, line: str, object_id: str, radius: bool = True) -> Any:
    x, y, w, h = box
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape._element.nvSpPr.cNvPr.set("name", object_id)
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(fill, "FFFFFF")
    shape.line.color.rgb = hex_to_rgb(line, "D7E2F0")
    shape.line.width = Pt(0.8)
    return shape


def add_title(slide: Any, row: dict[str, Any], profile: dict[str, Any], colors: dict[str, str]) -> None:
    add_text(slide, str(row.get("title", "")), (0.65, 0.34, 12.0, 0.7), "title", profile, colors["text"], f"s{row['slide_number']:03d}_title", vertical="top")
    if row.get("key_message"):
        add_text(slide, str(row["key_message"]), (0.68, 1.02, 11.7, 0.42), "subtitle", profile, colors["muted"], f"s{row['slide_number']:03d}_message", vertical="top")


def add_visual(slide: Any, slot: dict[str, Any], base_dir: Path, default_box: tuple[float, float, float, float], state: dict[str, int], object_name: str) -> None:
    asset = Path(str(slot["asset_path"]))
    asset_path = asset if asset.is_absolute() else base_dir / asset
    raw_box = slot.get("bbox_in") or list(default_box)
    box = tuple(float(value) for value in raw_box)
    fitted = fit_contain(asset_path, box)
    picture = slide.shapes.add_picture(str(asset_path), Inches(fitted[0]), Inches(fitted[1]), width=Inches(fitted[2]), height=Inches(fitted[3]))
    set_object_name(picture, object_name)
    state["image_objects"] += 1
    if slot.get("backend") == "image-2":
        state["image2_assets"] += 1


def validate_visual_slot(slot: dict[str, Any], base_dir: Path, policy: str, slide_number: int) -> list[str]:
    errors: list[str] = []
    backend = str(slot.get("backend") or "provided")
    source_type = str(slot.get("source_type") or "provided_asset")
    status = str(slot.get("status") or "planned")
    if policy == "native-only" and backend == "image-2":
        errors.append(f"slide_{slide_number}:native_only_rejects_image2")
    if backend == "image-2":
        if source_type != "imagegen_asset":
            errors.append(f"slide_{slide_number}:image2_source_type_must_be_imagegen_asset")
        if status not in {"generated", "validated"}:
            errors.append(f"slide_{slide_number}:image2_asset_not_generated")
        prompt = slot.get("prompt_record")
        if not prompt:
            errors.append(f"slide_{slide_number}:image2_prompt_record_missing")
        else:
            prompt_path = Path(str(prompt))
            prompt_path = prompt_path if prompt_path.is_absolute() else base_dir / prompt_path
            if not prompt_path.is_file():
                errors.append(f"slide_{slide_number}:image2_prompt_record_not_found")
    asset = slot.get("asset_path")
    if not asset:
        errors.append(f"slide_{slide_number}:visual_asset_path_missing")
    else:
        asset_path = Path(str(asset))
        asset_path = asset_path if asset_path.is_absolute() else base_dir / asset_path
        if not asset_path.is_file():
            errors.append(f"slide_{slide_number}:visual_asset_not_found")
    return errors


def build_cover(slide: Any, row: dict[str, Any], profile: dict[str, Any], colors: dict[str, str], base_dir: Path, state: dict[str, int]) -> None:
    visual = row.get("visual_slot")
    title_box = (0.85, 1.35, 6.8 if visual else 11.6, 1.45)
    add_text(slide, str(row.get("title", "")), title_box, "hero", profile, colors["text"], f"s{row['slide_number']:03d}_hero", vertical="top")
    if row.get("subtitle"):
        add_text(slide, str(row["subtitle"]), (0.9, 3.0, 6.4 if visual else 10.8, 0.8), "subtitle", profile, colors["muted"], f"s{row['slide_number']:03d}_subtitle", vertical="top")
    if row.get("meta"):
        add_text(slide, "  ·  ".join(str(value) for value in row["meta"]), (0.9, 6.45, 8.0, 0.35), "caption", profile, colors["muted"], f"s{row['slide_number']:03d}_meta")
    if visual:
        add_visual(slide, visual, base_dir, (8.0, 0.9, 4.5, 5.8), state, f"s{row['slide_number']:03d}_visual")


def build_content(slide: Any, row: dict[str, Any], profile: dict[str, Any], colors: dict[str, str], state: dict[str, int]) -> None:
    add_title(slide, row, profile, colors)
    cards = row.get("cards", [])
    columns = 2 if len(cards) <= 4 else 3
    rows = (len(cards) + columns - 1) // columns
    gap = 0.25
    card_w = (12.0 - gap * (columns - 1)) / columns
    card_h = (5.35 - gap * (rows - 1)) / max(1, rows)
    for index, card in enumerate(cards):
        col, row_index = index % columns, index // columns
        x = 0.65 + col * (card_w + gap)
        y = 1.62 + row_index * (card_h + gap)
        add_box(slide, (x, y, card_w, card_h), colors["panel"], colors["primary"], f"s{row['slide_number']:03d}_card_{index+1}")
        add_text(slide, str(card.get("title", "")), (x + 0.22, y + 0.18, card_w - 0.44, 0.45), "header", profile, colors["text"], f"s{row['slide_number']:03d}_card_{index+1}_title", vertical="top")
        add_text(slide, str(card.get("body", "")), (x + 0.22, y + 0.75, card_w - 0.44, card_h - 0.95), "body", profile, colors["muted"], f"s{row['slide_number']:03d}_card_{index+1}_body", vertical="top")
        state["native_shapes"] += 1


def build_process(slide: Any, row: dict[str, Any], profile: dict[str, Any], colors: dict[str, str], state: dict[str, int]) -> None:
    add_title(slide, row, profile, colors)
    steps = row.get("steps", [])
    gap = 0.28
    width = (12.0 - gap * (len(steps) - 1)) / max(1, len(steps))
    for index, step in enumerate(steps):
        x = 0.65 + index * (width + gap)
        add_box(slide, (x, 2.0, width, 3.65), colors["panel"], colors["primary"], f"s{row['slide_number']:03d}_step_{index+1}")
        add_text(slide, f"{index + 1:02d}", (x + 0.18, 2.2, 0.7, 0.42), "label", profile, colors["primary"], f"s{row['slide_number']:03d}_step_{index+1}_number")
        add_text(slide, str(step.get("title", "")), (x + 0.18, 2.85, width - 0.36, 0.55), "header", profile, colors["text"], f"s{row['slide_number']:03d}_step_{index+1}_title", vertical="top")
        add_text(slide, str(step.get("body", "")), (x + 0.18, 3.55, width - 0.36, 1.55), "body", profile, colors["muted"], f"s{row['slide_number']:03d}_step_{index+1}_body", vertical="top")
        state["native_shapes"] += 1
        if index < len(steps) - 1:
            connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + width), Inches(3.8), Inches(x + width + gap), Inches(3.8))
            set_object_name(connector, f"s{row['slide_number']:03d}_connector_{index+1}")
            connector.line.color.rgb = hex_to_rgb(colors["secondary"], "60A5FA")
            connector.line.width = Pt(1.5)
            state["connectors"] += 1


def add_bullet_list(slide: Any, values: list[Any], box: tuple[float, float, float, float], profile: dict[str, Any], color: str, object_id: str) -> None:
    add_text(slide, "\n".join(f"• {value}" for value in values), box, "body", profile, color, object_id, vertical="top")


def build_comparison(slide: Any, row: dict[str, Any], profile: dict[str, Any], colors: dict[str, str], state: dict[str, int]) -> None:
    add_title(slide, row, profile, colors)
    sides = [row.get("left", {}), row.get("right", {})]
    for index, side in enumerate(sides):
        x = 0.65 + index * 6.15
        add_box(slide, (x, 1.75, 5.85, 4.9), colors["panel"], colors["primary"] if index else colors["secondary"], f"s{row['slide_number']:03d}_compare_{index+1}")
        add_text(slide, str(side.get("title", "")), (x + 0.3, 2.05, 5.25, 0.6), "header", profile, colors["text"], f"s{row['slide_number']:03d}_compare_{index+1}_title")
        add_bullet_list(slide, list(side.get("points", [])), (x + 0.35, 2.85, 5.1, 3.35), profile, colors["muted"], f"s{row['slide_number']:03d}_compare_{index+1}_points")
        state["native_shapes"] += 1


def add_native_chart(slide: Any, chart_spec: dict[str, Any], box: tuple[float, float, float, float], colors: dict[str, str], state: dict[str, int], object_name: str) -> None:
    chart_data = ChartData()
    chart_data.categories = [str(value) for value in chart_spec.get("categories", [])]
    for series in chart_spec.get("series", []):
        chart_data.add_series(str(series.get("name", "Series")), [float(value) for value in series.get("values", [])])
    chart_type = {
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "line": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
    }.get(str(chart_spec.get("type", "column")), XL_CHART_TYPE.COLUMN_CLUSTERED)
    x, y, w, h = box
    chart_shape = slide.shapes.add_chart(chart_type, Inches(x), Inches(y), Inches(w), Inches(h), chart_data)
    set_object_name(chart_shape, object_name)
    chart = chart_shape.chart
    chart.has_legend = len(chart_spec.get("series", [])) > 1 or chart_type == XL_CHART_TYPE.PIE
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    chart.has_title = bool(chart_spec.get("title"))
    if chart.has_title:
        chart.chart_title.text_frame.text = str(chart_spec["title"])
    state["office_charts"] += 1


def build_data(slide: Any, row: dict[str, Any], profile: dict[str, Any], colors: dict[str, str], state: dict[str, int]) -> None:
    add_title(slide, row, profile, colors)
    metrics = row.get("metrics", [])
    width = 12.0 / max(1, len(metrics))
    for index, metric in enumerate(metrics):
        x = 0.65 + index * width
        add_text(slide, str(metric.get("value", "")), (x, 1.55, width - 0.2, 0.72), "section_title", profile, colors["primary"], f"s{row['slide_number']:03d}_metric_{index+1}_value")
        add_text(slide, str(metric.get("label", "")), (x, 2.25, width - 0.2, 0.38), "label", profile, colors["muted"], f"s{row['slide_number']:03d}_metric_{index+1}_label", vertical="top")
    if row.get("chart"):
        add_native_chart(slide, row["chart"], (0.8, 3.0, 11.7, 3.65), colors, state, f"s{row['slide_number']:03d}_chart")


def build_table(slide: Any, row: dict[str, Any], profile: dict[str, Any], table_policy: dict[str, Any], colors: dict[str, str], state: dict[str, int]) -> None:
    add_title(slide, row, profile, colors)
    headers = [str(value) for value in row.get("headers", [])]
    data_rows = [[str(value) for value in values] for values in row.get("rows", [])]
    columns = len(headers)
    table_shape = slide.shapes.add_table(len(data_rows) + 1, columns, Inches(0.7), Inches(1.65), Inches(11.95), Inches(4.95))
    table_shape._element.nvGraphicFramePr.cNvPr.set("name", f"s{row['slide_number']:03d}_table")
    table = table_shape.table
    column_types = list(row.get("column_types", []))
    size, east_asia, latin, paragraph_policy = role_settings(profile, "table")
    for row_index, values in enumerate([headers, *data_rows]):
        for col_index, value in enumerate(values):
            cell = table.cell(row_index, col_index)
            cell.text = value
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            cell.fill.fore_color.rgb = hex_to_rgb(colors["primary"] if row_index == 0 else colors["panel"], "FFFFFF")
            paragraph = cell.text_frame.paragraphs[0]
            kind = column_types[col_index] if col_index < len(column_types) else ("index" if col_index == 0 else "text")
            expected = table_policy["header_alignment"] if row_index == 0 else table_policy["numeric_alignment"] if kind == "numeric" else table_policy["index_alignment"] if kind == "index" else table_policy["text_alignment"]
            paragraph.alignment = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}[expected]
            paragraph.line_spacing = float(paragraph_policy.get("line_spacing_multiple", 1.0))
            paragraph.space_before = Pt(0)
            paragraph.space_after = Pt(0)
            for run in paragraph.runs:
                run.font.name = latin
                run.font.size = Pt(size)
                run.font.bold = row_index == 0
                run.font.color.rgb = hex_to_rgb("FFFFFF" if row_index == 0 else colors["text"], "111827")
                set_run_typefaces(run, east_asia, latin, profile["fonts"].get("complex_script") or latin)
    state["native_tables"] += 1


def build_architecture(slide: Any, row: dict[str, Any], profile: dict[str, Any], colors: dict[str, str], state: dict[str, int]) -> None:
    add_title(slide, row, profile, colors)
    layers = row.get("layers", [])
    height = 4.9 / max(1, len(layers))
    for index, layer in enumerate(layers):
        y = 1.65 + index * height
        add_box(slide, (1.0, y, 11.3, height - 0.16), colors["panel"], colors["primary"], f"s{row['slide_number']:03d}_layer_{index+1}", radius=False)
        add_text(slide, str(layer.get("title", "")), (1.25, y + 0.14, 2.3, height - 0.42), "header", profile, colors["primary"], f"s{row['slide_number']:03d}_layer_{index+1}_title")
        add_text(slide, str(layer.get("body", "")), (3.6, y + 0.14, 8.35, height - 0.42), "body", profile, colors["text"], f"s{row['slide_number']:03d}_layer_{index+1}_body")
        state["native_shapes"] += 1


def build_closing(slide: Any, row: dict[str, Any], profile: dict[str, Any], colors: dict[str, str], state: dict[str, int]) -> None:
    add_text(slide, str(row.get("title", "")), (1.0, 1.2, 11.3, 1.2), "hero", profile, colors["text"], f"s{row['slide_number']:03d}_closing_title", align="center")
    if row.get("key_message"):
        add_text(slide, str(row["key_message"]), (1.4, 2.55, 10.5, 0.75), "subtitle", profile, colors["muted"], f"s{row['slide_number']:03d}_closing_message", align="center")
    actions = row.get("actions", [])
    width = 9.6 / max(1, len(actions))
    for index, action in enumerate(actions):
        x = 1.85 + index * width
        add_box(slide, (x, 4.05, width - 0.25, 1.25), colors["panel"], colors["primary"], f"s{row['slide_number']:03d}_action_{index+1}")
        add_text(slide, str(action), (x + 0.14, 4.25, width - 0.53, 0.75), "body", profile, colors["text"], f"s{row['slide_number']:03d}_action_{index+1}_text", align="center")
        state["native_shapes"] += 1


def apply_text_override(shape: Any, element: dict[str, Any], profile: dict[str, Any]) -> None:
    if not getattr(shape, "has_text_frame", False):
        return
    role = str(element.get("role") or object_role(shape) or "body")
    if role not in ROLE_SIZE_KEY:
        role = "body"
    default_size, east_asia, latin, paragraph_policy = role_settings(profile, role)
    style = element.get("style") if isinstance(element.get("style"), dict) else {}
    size = float(style.get("font_size_pt", default_size))
    font_family = str(style.get("font_family") or latin)
    font_color = str(style.get("font_color") or "111827")
    bold = bool(style.get("bold", role in {"hero", "section_title", "title", "header"}))
    align = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}.get(str(style.get("align", "left")), PP_ALIGN.LEFT)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    vertical = str(style.get("vertical_align", "middle"))
    frame.vertical_anchor = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}.get(vertical, MSO_ANCHOR.MIDDLE)
    for index, line in enumerate(str(element.get("text", "")).split("\n")):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        run = paragraph.add_run(); run.text = line
        run.font.name = font_family
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = hex_to_rgb(font_color, "111827")
        set_run_typefaces(run, east_asia, font_family, profile["fonts"].get("complex_script") or font_family)
        paragraph.alignment = align
        paragraph.line_spacing = float(style.get("line_spacing", paragraph_policy.get("line_spacing_multiple", 1.0)))
        paragraph.space_before = Pt(size * float(paragraph_policy.get("space_before_lines", 0)))
        paragraph.space_after = Pt(size * float(paragraph_policy.get("space_after_lines", 0)))


def materialize_scene_element(slide: Any, element: dict[str, Any], canvas_w: float, canvas_h: float) -> Any | None:
    """Create native PowerPoint objects that exist on the canvas but not in the base archetype."""
    element_id = str(element.get("id") or "")
    bbox = element.get("bbox")
    if not element_id or not isinstance(bbox, list) or len(bbox) != 4:
        return None
    x, y, w, h = (float(value) for value in bbox)
    left, top = Inches(x / canvas_w * SLIDE_W), Inches(y / canvas_h * SLIDE_H)
    width, height = Inches(max(1.0, w) / canvas_w * SLIDE_W), Inches(max(1.0, h) / canvas_h * SLIDE_H)
    kind = str(element.get("type") or "")
    if kind == "text":
        shape = slide.shapes.add_textbox(left, top, width, height)
        set_role(shape, str(element.get("role") or "body"), element_id)
        return shape
    if kind == "shape":
        geometry = str(element.get("geometry") or "rectangle").lower()
        shape_type = {
            "ellipse": MSO_SHAPE.OVAL,
            "oval": MSO_SHAPE.OVAL,
            "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
            "roundrect": MSO_SHAPE.ROUNDED_RECTANGLE,
        }.get(geometry, MSO_SHAPE.RECTANGLE)
        shape = slide.shapes.add_shape(shape_type, left, top, width, height)
        set_object_name(shape, element_id)
        return shape
    if kind == "connector":
        shape = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT,
            left,
            top,
            Inches((x + w) / canvas_w * SLIDE_W),
            Inches((y + h) / canvas_h * SLIDE_H),
        )
        set_object_name(shape, element_id)
        return shape
    return None


def apply_scene_overrides(slide: Any, scene: dict[str, Any], profile: dict[str, Any]) -> tuple[int, int, int]:
    canvas = scene.get("canvas", {})
    canvas_w = float(canvas.get("width", 1920))
    canvas_h = float(canvas.get("height", 1080))
    shapes = {object_id(shape): shape for shape in slide.shapes if object_id(shape)}
    applied = 0
    materialized_shapes = 0
    materialized_connectors = 0
    z_rows: list[tuple[int, Any]] = []
    for element in scene.get("elements", []):
        shape = shapes.get(str(element.get("id")))
        if shape is None:
            shape = materialize_scene_element(slide, element, canvas_w, canvas_h)
            if shape is None:
                continue
            shapes[str(element.get("id"))] = shape
            if str(element.get("type")) == "connector":
                materialized_connectors += 1
            else:
                materialized_shapes += 1
        bbox = element.get("bbox")
        if isinstance(bbox, list) and len(bbox) == 4:
            x, y, w, h = (float(value) for value in bbox)
            shape.left = Inches(x / canvas_w * SLIDE_W)
            shape.top = Inches(y / canvas_h * SLIDE_H)
            shape.width = Inches(w / canvas_w * SLIDE_W)
            shape.height = Inches(h / canvas_h * SLIDE_H)
        if "rotation" in element and hasattr(shape, "rotation"):
            shape.rotation = float(element["rotation"])
        style = element.get("style") if isinstance(element.get("style"), dict) else {}
        if style.get("fill") and hasattr(shape, "fill"):
            shape.fill.solid(); shape.fill.fore_color.rgb = hex_to_rgb(str(style["fill"]), "FFFFFF")
        if style.get("line") and hasattr(shape, "line"):
            shape.line.color.rgb = hex_to_rgb(str(style["line"]), "D7E2F0")
        if style.get("line_width_pt") is not None and hasattr(shape, "line"):
            shape.line.width = Pt(float(style["line_width_pt"]))
        if "text" in element:
            apply_text_override(shape, element, profile)
        z_rows.append((int(element.get("z_index", 0)), shape))
        applied += 1
    tree = slide.shapes._spTree
    for _, shape in sorted(z_rows, key=lambda row: row[0]):
        tree.remove(shape._element)
        tree.append(shape._element)
    return applied, materialized_shapes, materialized_connectors


def validate_spec(spec: dict[str, Any], base_dir: Path) -> list[str]:
    errors: list[str] = []
    policy = str(spec.get("visual_asset_policy") or "native-image-assisted")
    if policy not in POLICIES:
        errors.append(f"invalid_visual_asset_policy:{policy}")
    try:
        resolve_style(str(spec.get("style_id") or "consulting-blue-white"), spec.get("style_variant"))
    except ValueError as exc:
        errors.append(str(exc))
    slides = spec.get("slides")
    if not isinstance(slides, list) or not slides:
        return [*errors, "slides_missing"]
    image2_count = 0
    for expected_number, row in enumerate(slides, start=1):
        number = int(row.get("slide_number", 0))
        if number != expected_number:
            errors.append(f"slide_number_not_contiguous:{number}:{expected_number}")
        archetype = str(row.get("archetype") or "")
        if archetype not in ARCHETYPES:
            errors.append(f"slide_{number}:unsupported_archetype:{archetype}")
        if not str(row.get("title") or "").strip():
            errors.append(f"slide_{number}:title_missing")
        try:
            resolve_style(str(row.get("style_id") or spec.get("style_id") or "consulting-blue-white"), row.get("style_variant") or spec.get("style_variant"))
        except ValueError as exc:
            errors.append(f"slide_{number}:{exc}")
        visual = row.get("visual_slot")
        if visual:
            errors.extend(validate_visual_slot(visual, base_dir, policy, number))
            if visual.get("backend") == "image-2":
                image2_count += 1
        if archetype == "content-structured" and not 2 <= len(row.get("cards", [])) <= 6:
            errors.append(f"slide_{number}:cards_must_be_2_to_6")
        if archetype == "process-flow" and not 3 <= len(row.get("steps", [])) <= 5:
            errors.append(f"slide_{number}:steps_must_be_3_to_5")
        if archetype == "data-callouts" and not 2 <= len(row.get("metrics", [])) <= 5:
            errors.append(f"slide_{number}:metrics_must_be_2_to_5")
        if archetype == "table" and (not row.get("headers") or not row.get("rows")):
            errors.append(f"slide_{number}:table_data_missing")
    if policy == "image-led-editable" and image2_count == 0:
        errors.append("image_led_editable_requires_image2_asset")
    return errors


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--spec", required=True)
    parser.add_argument("--out-pptx", required=True)
    parser.add_argument("--base-dir")
    parser.add_argument("--scene-dir", help="Optional scene directory whose element geometry/text/style overrides the generated layout.")
    parser.add_argument("--session", help="Optional session to mark editor canvas rebuilds complete.")
    parser.add_argument("--report")
    parser.add_argument("--design-report", help="Optional Anti-AI-slop design-gate report path.")
    args = parser.parse_args()

    spec_path = Path(args.spec).resolve()
    base_dir = Path(args.base_dir).resolve() if args.base_dir else spec_path.parent
    spec = read_json(spec_path)
    if args.session:
        session_path = Path(args.session).resolve()
        metadata_path = session_path / "metadata.json"
        if metadata_path.is_file():
            editor_metadata = read_json(metadata_path)
            out_candidate = Path(args.out_pptx).resolve()
            final_dir = (session_path / "final").resolve()
            if (
                editor_metadata.get("editor_workflow_mode") == "canvas-first"
                and (out_candidate == final_dir or final_dir in out_candidate.parents)
                and editor_metadata.get("editor_export_approval") != "approved"
            ):
                report = {"schema_version": 1, "status": "FAIL", "errors": ["canvas_export:not_approved"]}
                print(json.dumps(report, ensure_ascii=False))
                raise SystemExit(2)
    errors = validate_spec(spec, base_dir)
    if errors:
        report = {"schema_version": 1, "status": "FAIL", "spec": str(spec_path), "errors": errors}
        if args.report:
            out_report = Path(args.report).resolve()
            out_report.parent.mkdir(parents=True, exist_ok=True)
            out_report.write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
        print(json.dumps(report, ensure_ascii=False))
        raise SystemExit(2)

    scene_dir = Path(args.scene_dir).resolve() if args.scene_dir else None
    scenes = [read_json(path) for path in sorted(scene_dir.glob("*.scene.json"))] if scene_dir and scene_dir.is_dir() else []
    design_report = analyze_design(spec, scenes)
    design_report_path = Path(args.design_report).resolve() if args.design_report else (Path(args.session).resolve() / "reports" / "design-quality.json" if args.session else None)
    if design_report_path:
        design_report_path.parent.mkdir(parents=True, exist_ok=True)
        design_report_path.write_text(json.dumps(design_report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    if design_report["status"] == "FAIL":
        print(json.dumps(design_report, ensure_ascii=False))
        raise SystemExit(2)

    style_id = str(spec.get("style_id") or "consulting-blue-white")
    colors, profile, table_policy, profile_id, table_id, style_variant = load_design(style_id, spec.get("style_variant"), spec.get("typography_profile"), spec.get("table_profile"))
    presentation = Presentation()
    presentation.slide_width = Inches(SLIDE_W)
    presentation.slide_height = Inches(SLIDE_H)
    blank = presentation.slide_layouts[6]
    state: dict[str, int] = {"native_shapes": 0, "connectors": 0, "native_tables": 0, "office_charts": 0, "image_objects": 0, "image2_assets": 0, "canvas_overrides_applied": 0}
    for row in spec["slides"]:
        slide_style_id = str(row.get("style_id") or style_id)
        slide_colors, slide_profile, slide_table_policy, _, _, _ = load_design(
            slide_style_id,
            row.get("style_variant") or spec.get("style_variant"),
            spec.get("typography_profile"),
            spec.get("table_profile"),
        )
        builders = {
            "cover": lambda slide, item: build_cover(slide, item, slide_profile, slide_colors, base_dir, state),
            "content-structured": lambda slide, item: build_content(slide, item, slide_profile, slide_colors, state),
            "process-flow": lambda slide, item: build_process(slide, item, slide_profile, slide_colors, state),
            "comparison-two-zone": lambda slide, item: build_comparison(slide, item, slide_profile, slide_colors, state),
            "data-callouts": lambda slide, item: build_data(slide, item, slide_profile, slide_colors, state),
            "table": lambda slide, item: build_table(slide, item, slide_profile, slide_table_policy, slide_colors, state),
            "architecture": lambda slide, item: build_architecture(slide, item, slide_profile, slide_colors, state),
            "closing-action": lambda slide, item: build_closing(slide, item, slide_profile, slide_colors, state),
        }
        slide = presentation.slides.add_slide(blank)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = hex_to_rgb(slide_colors["background"], "FFFFFF")
        builders[str(row["archetype"])](slide, row)
        if row.get("visual_slot") and row["archetype"] != "cover":
            add_visual(slide, row["visual_slot"], base_dir, (9.4, 5.2, 3.1, 1.75), state, f"s{row['slide_number']:03d}_visual")
        if scene_dir:
            scene_path = scene_dir / f"slide-{int(row['slide_number']):03d}.scene.json"
            if scene_path.is_file():
                applied, added_shapes, added_connectors = apply_scene_overrides(slide, read_json(scene_path), profile)
                state["canvas_overrides_applied"] += applied
                state["native_shapes"] += added_shapes
                state["connectors"] += added_connectors

    out_pptx = Path(args.out_pptx).resolve()
    out_pptx.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(out_pptx)
    patch_theme_east_asian_fonts(out_pptx, profile["fonts"]["heading_east_asia"], profile["fonts"]["body_east_asia"])
    if args.session:
        session = Path(args.session).resolve()
        metadata_path = session / "metadata.json"
        metadata = read_json(metadata_path)
        try:
            artifact_ref = str(out_pptx.relative_to(session))
        except ValueError:
            artifact_ref = str(out_pptx)
        built_at = datetime.now(timezone.utc).isoformat()
        variant = metadata.setdefault("variants", {}).setdefault("editable", {"status": "not_built", "artifact": None, "slides": {}})
        variant["status"] = "built"
        variant["artifact"] = artifact_ref
        for row in spec["slides"]:
            number = int(row["slide_number"])
            scene_path = session / "scenes" / f"slide-{number:03d}.scene.json"
            scene_hash = read_json(scene_path).get("dependencies", {}).get("scene_hash") if scene_path.is_file() else None
            variant.setdefault("slides", {})[str(number)] = {
                "status": "built", "scene_hash": scene_hash, "artifact": artifact_ref,
                "artifact_hash": sha256(out_pptx), "built_at": built_at,
            }
        if isinstance(metadata.get("editor"), dict):
            metadata["editor"]["status"] = "ready"
        metadata_path.write_text(json.dumps(metadata, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
        events_path = session / "reports" / "editor-events.jsonl"
        events_path.parent.mkdir(parents=True, exist_ok=True)
        with events_path.open("a", encoding="utf-8") as handle:
            handle.write(json.dumps({
                "event": "editor.rebuild.completed", "artifact": artifact_ref,
                "slides": [int(row["slide_number"]) for row in spec["slides"]], "created_at": built_at,
            }, ensure_ascii=False) + "\n")
    state["native_text_objects"] = sum(
        1
        for slide in presentation.slides
        for shape in slide.shapes
        if getattr(shape, "has_text_frame", False) and shape.text_frame.text.strip()
    )
    report = {
        "schema_version": 1,
        "status": "PASS",
        "route": "native-editable-deck",
        "pptx": str(out_pptx),
        "slide_count": len(spec["slides"]),
        "style_id": style_id,
        "style_variant": style_variant,
        "typography_profile": profile_id,
        "table_profile": table_id,
        "visual_asset_policy": spec.get("visual_asset_policy") or "native-image-assisted",
        "design_quality": design_report["status"],
        "session": str(Path(args.session).resolve()) if args.session else None,
        **state,
        "errors": [],
    }
    if args.report:
        report_path = Path(args.report).resolve()
        report_path.parent.mkdir(parents=True, exist_ok=True)
        report_path.write_text(json.dumps(report, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")
    print(json.dumps(report, ensure_ascii=False))


if __name__ == "__main__":
    main()
