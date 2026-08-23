#!/usr/bin/env python3
"""Validate semantic typography tokens and table alignment in a PPTX."""

from __future__ import annotations

import argparse
import json
import math
import re
from collections import defaultdict
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_PROFILES = ROOT / "styles" / "typography-profiles.json"
DEFAULT_CATALOG = ROOT / "styles" / "catalog.json"
ROLE_MARKER = re.compile(r"\[pf-role=([a-z0-9_-]+)]", re.I)
NUMERIC = re.compile(r"^\s*[¥￥$€£]?\(?[-+]?\d[\d,]*(?:\.\d+)?%?\)?\s*$")
ROLE_SIZE_KEY = {
    "hero": "hero",
    "section_title": "section_title",
    "section-title": "section_title",
    "title": "page_title",
    "page_title": "page_title",
    "page-title": "page_title",
    "subtitle": "subtitle",
    "header": "minor_title",
    "minor_title": "minor_title",
    "minor-title": "minor_title",
    "body": "body",
    "label": "label",
    "micro_label": "caption",
    "micro-label": "caption",
    "caption": "caption",
    "table": "table",
}


def read_json(path: Path) -> dict[str, Any]:
    value = json.loads(path.read_text(encoding="utf-8"))
    if not isinstance(value, dict):
        raise ValueError(f"expected JSON object: {path}")
    return value


def select_policies(args: argparse.Namespace) -> tuple[dict[str, Any], dict[str, Any], str, str]:
    profiles = read_json(Path(args.profiles))
    typography = {row["id"]: row for row in profiles.get("typography_profiles", [])}
    tables = {row["id"]: row for row in profiles.get("table_profiles", [])}
    profile_id = args.profile
    table_id = args.table_profile
    if args.style_id:
        catalog = read_json(Path(args.catalog))
        style = next((row for row in catalog.get("styles", []) if row.get("id") == args.style_id), None)
        if not style:
            raise ValueError(f"unknown style_id: {args.style_id}")
        profile_id = profile_id or style.get("typography_profile")
        table_id = table_id or style.get("table_profile")
    if not profile_id or profile_id not in typography:
        raise ValueError(f"unknown or missing typography profile: {profile_id!r}")
    if not table_id or table_id not in tables:
        raise ValueError(f"unknown or missing table profile: {table_id!r}")
    return typography[profile_id], tables[table_id], profile_id, table_id


def add_finding(store: dict[tuple[str, str], dict[str, Any]], code: str, severity: str, sample: dict[str, Any]) -> None:
    key = (code, severity)
    finding = store.setdefault(key, {"code": code, "severity": severity, "count": 0, "samples": []})
    finding["count"] += 1
    if len(finding["samples"]) < 8:
        finding["samples"].append(sample)


def role_from_shape(shape: Any) -> str | None:
    name = str(getattr(shape, "name", ""))
    marker = ROLE_MARKER.search(name)
    if marker:
        return marker.group(1).lower()
    lowered = name.lower()
    for token, role in (("subtitle", "subtitle"), ("caption", "caption"), ("label", "label"), ("title", "title"), ("body", "body")):
        if token in lowered:
            return role
    if getattr(shape, "is_placeholder", False):
        placeholder = str(shape.placeholder_format.type).upper()
        if "SUBTITLE" in placeholder:
            return "subtitle"
        if "TITLE" in placeholder:
            return "title"
        if "BODY" in placeholder or "OBJECT" in placeholder or "TEXT" in placeholder:
            return "body"
    return None


def alignment_name(value: Any) -> str | None:
    return {
        PP_ALIGN.LEFT: "left",
        PP_ALIGN.CENTER: "center",
        PP_ALIGN.RIGHT: "right",
        PP_ALIGN.JUSTIFY: "justify",
        PP_ALIGN.DISTRIBUTE: "distribute",
    }.get(value)


def check_run_sizes(
    paragraphs: Any,
    role: str,
    profile: dict[str, Any],
    location: dict[str, Any],
    findings: dict[tuple[str, str], dict[str, Any]],
    counts: dict[str, int],
) -> None:
    size_key = ROLE_SIZE_KEY.get(role)
    if not size_key:
        add_finding(findings, "text_role_unresolved", "not_checked", {**location, "role": role})
        return
    expected = float(profile["tokens"][size_key])
    grid = float(profile.get("font_size_grid", 0.5))
    for paragraph_index, paragraph in enumerate(paragraphs):
        if paragraph.text and not paragraph.runs:
            add_finding(findings, "font_size_unresolved", "not_checked", {**location, "paragraph": paragraph_index, "role": role})
        for run_index, run in enumerate(paragraph.runs):
            if not run.text.strip():
                continue
            counts["text_runs"] += 1
            if run.font.size is None:
                add_finding(findings, "font_size_unresolved", "not_checked", {**location, "paragraph": paragraph_index, "run": run_index, "role": role})
                continue
            size = float(run.font.size.pt)
            grid_units = size / grid
            if not math.isclose(grid_units, round(grid_units), abs_tol=0.04):
                add_finding(findings, "font_size_off_grid", "warning", {**location, "role": role, "font_size": round(size, 3), "grid": grid})
            if size + 0.04 < expected:
                add_finding(findings, "font_size_below_token", "warning", {**location, "role": role, "font_size": round(size, 3), "size_key": size_key, "expected": expected})


def check_paragraph_policy(
    paragraphs: Any,
    role: str,
    profile: dict[str, Any],
    location: dict[str, Any],
    findings: dict[tuple[str, str], dict[str, Any]],
) -> None:
    group = "title" if ROLE_SIZE_KEY.get(role) in {"hero", "section_title", "page_title", "subtitle", "minor_title"} else "caption" if ROLE_SIZE_KEY.get(role) in {"label", "caption"} else "body"
    expected = profile.get("paragraph", {}).get(group, {})
    expected_spacing = expected.get("line_spacing_multiple")
    for paragraph_index, paragraph in enumerate(paragraphs):
        spacing = paragraph.line_spacing
        if spacing is None:
            add_finding(findings, "line_spacing_unresolved", "not_checked", {**location, "paragraph": paragraph_index, "role": role})
        elif isinstance(spacing, float) and expected_spacing is not None and not math.isclose(spacing, float(expected_spacing), abs_tol=0.03):
            add_finding(findings, "line_spacing_mismatch", "warning", {**location, "paragraph": paragraph_index, "role": role, "actual": round(spacing, 3), "expected": expected_spacing})


def expected_cell_alignment(row: int, col: int, text: str, table_policy: dict[str, Any]) -> str:
    if row == 0:
        return str(table_policy["header_alignment"])
    if NUMERIC.fullmatch(text or ""):
        return str(table_policy["numeric_alignment"])
    if col == 0:
        return str(table_policy["index_alignment"])
    return str(table_policy["text_alignment"])


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("pptx")
    parser.add_argument("--profile")
    parser.add_argument("--table-profile")
    parser.add_argument("--style-id")
    parser.add_argument("--profiles", default=str(DEFAULT_PROFILES))
    parser.add_argument("--catalog", default=str(DEFAULT_CATALOG))
    parser.add_argument("--out")
    parser.add_argument("--fail-on-warning", action="store_true")
    args = parser.parse_args()

    pptx_path = Path(args.pptx).resolve()
    findings: dict[tuple[str, str], dict[str, Any]] = {}
    counts: dict[str, int] = defaultdict(int)
    try:
        profile, table_policy, profile_id, table_id = select_policies(args)
        presentation = Presentation(pptx_path)
        for slide_index, slide in enumerate(presentation.slides, start=1):
            for shape in slide.shapes:
                location = {"slide": slide_index, "shape": getattr(shape, "name", "")}
                if getattr(shape, "has_table", False):
                    table = shape.table
                    for row_index, row in enumerate(table.rows):
                        for col_index, cell in enumerate(row.cells):
                            counts["table_cells"] += 1
                            cell_location = {**location, "row": row_index, "column": col_index}
                            if cell.vertical_anchor is None:
                                add_finding(findings, "table_vertical_alignment_unresolved", "not_checked", cell_location)
                            elif cell.vertical_anchor != MSO_ANCHOR.MIDDLE:
                                add_finding(findings, "table_vertical_alignment_mismatch", "warning", {**cell_location, "expected": "middle"})
                            expected_alignment = expected_cell_alignment(row_index, col_index, cell.text, table_policy)
                            for paragraph_index, paragraph in enumerate(cell.text_frame.paragraphs):
                                actual_alignment = alignment_name(paragraph.alignment)
                                if actual_alignment is None:
                                    add_finding(findings, "table_horizontal_alignment_unresolved", "not_checked", {**cell_location, "paragraph": paragraph_index, "expected": expected_alignment})
                                elif actual_alignment != expected_alignment:
                                    add_finding(findings, "table_horizontal_alignment_mismatch", "warning", {**cell_location, "paragraph": paragraph_index, "actual": actual_alignment, "expected": expected_alignment})
                                if paragraph.level:
                                    add_finding(findings, "table_paragraph_special_indent", "warning", {**cell_location, "paragraph": paragraph_index, "level": paragraph.level})
                            check_run_sizes(cell.text_frame.paragraphs, "table", profile, cell_location, findings, counts)
                    continue
                if not getattr(shape, "has_text_frame", False) or not shape.text_frame.text.strip():
                    continue
                role = role_from_shape(shape)
                if not role:
                    add_finding(findings, "text_role_unresolved", "not_checked", location)
                    continue
                check_run_sizes(shape.text_frame.paragraphs, role, profile, location, findings, counts)
                check_paragraph_policy(shape.text_frame.paragraphs, role, profile, location, findings)
        if not counts["text_runs"] and not counts["table_cells"]:
            add_finding(findings, "no_typography_content", "not_checked", {"pptx": str(pptx_path)})
        findings_list = sorted(findings.values(), key=lambda row: (row["severity"], row["code"]))
        warning_count = sum(row["count"] for row in findings_list if row["severity"] == "warning")
        not_checked_count = sum(row["count"] for row in findings_list if row["severity"] == "not_checked")
        status = "WARN" if warning_count else "NOT_CHECKED" if not_checked_count else "PASS"
        report = {
            "schema_version": 1,
            "status": status,
            "pptx": str(pptx_path),
            "typography_profile": profile_id,
            "table_profile": table_id,
            "checked": dict(counts),
            "warning_count": warning_count,
            "not_checked_count": not_checked_count,
            "findings": findings_list,
        }
    except Exception as exc:
        report = {"schema_version": 1, "status": "FAIL", "pptx": str(pptx_path), "errors": [str(exc)]}

    payload = json.dumps(report, ensure_ascii=False, indent=2) + "\n"
    if args.out:
        out = Path(args.out).resolve()
        out.parent.mkdir(parents=True, exist_ok=True)
        out.write_text(payload, encoding="utf-8")
    print(payload, end="")
    if report["status"] == "FAIL" or (args.fail_on_warning and report["status"] == "WARN"):
        raise SystemExit(2)


if __name__ == "__main__":
    main()
