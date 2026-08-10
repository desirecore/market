#!/usr/bin/env python3
"""Analyze, plan, check, apply, and validate conservative native PPTX fills."""

from __future__ import annotations

import argparse
import copy
import json
import posixpath
import re
import shutil
import tempfile
import zipfile
from pathlib import Path

from lxml import etree as ET
from pptx import Presentation


P = "http://schemas.openxmlformats.org/presentationml/2006/main"
A = "http://schemas.openxmlformats.org/drawingml/2006/main"
R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
REL = "http://schemas.openxmlformats.org/package/2006/relationships"
XML = "http://www.w3.org/XML/1998/namespace"
NS = {"p": P, "a": A, "r": R, "rel": REL}


def write_json(path: Path, value: object) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(value, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def role_for_shape(shape: object) -> str:
    if getattr(shape, "is_placeholder", False):
        try:
            name = shape.placeholder_format.type.name.lower()
        except (AttributeError, ValueError):
            name = "placeholder"
        if "title" in name:
            return "title"
        if "sub" in name:
            return "subtitle"
        if any(word in name for word in ("body", "object", "text")):
            return "body"
        return name
    name = str(getattr(shape, "name", "")).lower()
    return "title" if "title" in name else "text"


def font_size_pt(text_frame: object, default: float = 18.0) -> float:
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size:
                return float(run.font.size.pt)
    return default


def capacity_units(width: int, height: int, font_pt: float) -> int:
    width_pt = max(1.0, width / 12700)
    height_pt = max(1.0, height / 12700)
    chars_per_line = max(1.0, width_pt / max(font_pt * 0.95, 1.0))
    lines = max(1, int(height_pt / max(font_pt * 1.25, 1.0)))
    return max(4, int(chars_per_line * lines))


def visual_units(text: str) -> float:
    units = 0.0
    for char in text:
        if char.isspace():
            units += 0.25
        elif ord(char) >= 0x2E80:
            units += 1.0
        elif char.isalnum():
            units += 0.55
        else:
            units += 0.45
    return units


def analyze_pptx(source: Path) -> dict[str, object]:
    prs = Presentation(str(source))
    slides: list[dict[str, object]] = []
    for slide_number, slide in enumerate(prs.slides, start=1):
        slots: list[dict[str, object]] = []
        for shape in slide.shapes:
            geometry = {
                "left": int(shape.left), "top": int(shape.top),
                "width": int(shape.width), "height": int(shape.height),
            }
            if getattr(shape, "has_text_frame", False):
                size = font_size_pt(shape.text_frame)
                slots.append({
                    "slot_id": f"s{slide_number:03d}_sh{shape.shape_id}",
                    "kind": "text",
                    "shape_id": int(shape.shape_id),
                    "shape_name": shape.name,
                    "role": role_for_shape(shape),
                    "geometry": geometry,
                    "font_size_pt": size,
                    "paragraph_count": len(shape.text_frame.paragraphs),
                    "capacity_visual_units": capacity_units(shape.width, shape.height, size),
                    "old_text": shape.text,
                })
            if getattr(shape, "has_table", False):
                rows = len(shape.table.rows)
                cols = len(shape.table.columns)
                cell_width = max(1, int(shape.width / max(cols, 1)))
                cell_height = max(1, int(shape.height / max(rows, 1)))
                for row_index, row in enumerate(shape.table.rows):
                    for col_index, cell in enumerate(row.cells):
                        size = font_size_pt(cell.text_frame, 14.0)
                        slots.append({
                            "slot_id": f"s{slide_number:03d}_tbl{shape.shape_id}_r{row_index}_c{col_index}",
                            "kind": "table-cell",
                            "shape_id": int(shape.shape_id),
                            "shape_name": shape.name,
                            "role": "table-cell",
                            "row": row_index,
                            "col": col_index,
                            "geometry": geometry,
                            "font_size_pt": size,
                            "paragraph_count": len(cell.text_frame.paragraphs),
                            "capacity_visual_units": capacity_units(cell_width, cell_height, size),
                            "old_text": cell.text,
                        })
        summary = " | ".join(slot["old_text"] for slot in slots if slot["old_text"] and slot["role"] in {"title", "subtitle", "body"})
        slides.append({"slide_number": slide_number, "text_summary": summary[:500], "slots": slots})
    return {
        "schema": "native_template_library.v1",
        "source_pptx": str(source),
        "slide_size": {"width": int(prs.slide_width), "height": int(prs.slide_height)},
        "slide_count": len(slides),
        "slides": slides,
        "boundaries": {
            "repeat_source_slide": False,
            "replace_images": False,
            "edit_charts": False,
            "edit_smartart": False,
        },
    }


def slot_map(library: dict[str, object]) -> dict[str, dict[str, object]]:
    return {
        str(slot["slot_id"]): slot
        for slide in library.get("slides", [])
        for slot in slide.get("slots", [])
    }


def check_plan(library: dict[str, object], plan: dict[str, object]) -> dict[str, object]:
    errors: list[dict[str, object]] = []
    warnings: list[dict[str, object]] = []
    slots = slot_map(library)
    max_slide = int(library.get("slide_count", 0))
    seen: set[int] = set()
    planned = plan.get("slides")
    font_policy = plan.get("font_policy")
    if font_policy is not None:
        if not isinstance(font_policy, dict) or not str(font_policy.get("font_face", "")).strip():
            errors.append({"code": "font_policy_invalid", "message": "font_policy requires a non-empty font_face."})
        elif font_policy.get("scope", "replaced-text") not in {"replaced-text", "all-selected-text", "theme-and-replaced"}:
            errors.append({"code": "font_scope_invalid", "scope": font_policy.get("scope")})
    if not isinstance(planned, list) or not planned:
        errors.append({"code": "slides_missing", "message": "Plan must contain at least one slide."})
        planned = []
    for plan_index, slide in enumerate(planned, start=1):
        source_slide = int(slide.get("source_slide", 0))
        if source_slide < 1 or source_slide > max_slide:
            errors.append({"code": "source_slide_missing", "plan_slide": plan_index, "source_slide": source_slide})
            continue
        if source_slide in seen:
            errors.append({"code": "repeat_source_slide_unsupported", "plan_slide": plan_index, "source_slide": source_slide})
        seen.add(source_slide)
        if not slide.get("layout_rationale"):
            warnings.append({"code": "layout_rationale_missing", "plan_slide": plan_index})
        for replacement in slide.get("replacements", []):
            slot_id = str(replacement.get("slot_id", ""))
            slot = slots.get(slot_id)
            if not slot:
                errors.append({"code": "slot_missing", "plan_slide": plan_index, "slot_id": slot_id})
                continue
            expected_prefix = f"s{source_slide:03d}_"
            if not slot_id.startswith(expected_prefix):
                errors.append({"code": "slot_slide_mismatch", "plan_slide": plan_index, "slot_id": slot_id})
                continue
            text = str(replacement.get("text", ""))
            capacity = float(slot.get("capacity_visual_units", 0) or 0)
            used = visual_units(text)
            ratio = used / capacity if capacity else 0
            if ratio > 1.15:
                warnings.append({
                    "code": "text_capacity", "plan_slide": plan_index, "slot_id": slot_id,
                    "used_visual_units": round(used, 2), "capacity_visual_units": capacity,
                    "ratio": round(ratio, 2),
                })
    return {
        "schema": "native_template_check.v1",
        "status": "FAIL" if errors else "WARN" if warnings else "PASS",
        "error_count": len(errors),
        "warning_count": len(warnings),
        "errors": errors,
        "warnings": warnings,
    }


def presentation_slide_parts(files: dict[str, bytes]) -> tuple[ET.Element, list[ET.Element], list[str]]:
    presentation = ET.fromstring(files["ppt/presentation.xml"])
    rels = ET.fromstring(files["ppt/_rels/presentation.xml.rels"])
    targets = {
        rel.get("Id"): posixpath.normpath(posixpath.join("ppt", rel.get("Target", "")))
        for rel in rels.findall(f"{{{REL}}}Relationship")
    }
    slide_list = presentation.find(f"{{{P}}}sldIdLst")
    if slide_list is None:
        raise ValueError("presentation has no slide list")
    slide_ids = list(slide_list)
    parts = [targets.get(node.get(f"{{{R}}}id"), "") for node in slide_ids]
    return presentation, slide_ids, parts


def set_rpr_typefaces(r_pr: ET.Element, font_face: str) -> None:
    for child in list(r_pr):
        if ET.QName(child).namespace == A and ET.QName(child).localname in {"latin", "ea", "cs"}:
            r_pr.remove(child)
    late_children = {"sym", "hlinkClick", "hlinkMouseOver", "extLst"}
    insertion_index = len(r_pr)
    for index, child in enumerate(r_pr):
        if ET.QName(child).namespace == A and ET.QName(child).localname in late_children:
            insertion_index = index
            break
    for offset, leaf in enumerate(("latin", "ea", "cs")):
        node = ET.Element(f"{{{A}}}{leaf}")
        node.set("typeface", font_face)
        r_pr.insert(insertion_index + offset, node)


def replace_text_body(tx_body: ET.Element, text: str, font_face: str | None = None) -> None:
    paragraphs = tx_body.findall(f"{{{A}}}p")
    template = paragraphs[0] if paragraphs else ET.Element(f"{{{A}}}p")
    template_ppr = template.find(f"{{{A}}}pPr")
    first_run = template.find(f"{{{A}}}r")
    template_rpr = first_run.find(f"{{{A}}}rPr") if first_run is not None else None
    template_end = template.find(f"{{{A}}}endParaRPr")
    for paragraph in paragraphs:
        tx_body.remove(paragraph)
    lines = text.splitlines() or [""]
    for line in lines:
        paragraph = ET.Element(f"{{{A}}}p")
        if template_ppr is not None:
            paragraph.append(copy.deepcopy(template_ppr))
        run = ET.SubElement(paragraph, f"{{{A}}}r")
        run_rpr = copy.deepcopy(template_rpr) if template_rpr is not None else None
        if font_face:
            run_rpr = run_rpr if run_rpr is not None else ET.Element(f"{{{A}}}rPr")
            set_rpr_typefaces(run_rpr, font_face)
        if run_rpr is not None:
            run.append(run_rpr)
        text_node = ET.SubElement(run, f"{{{A}}}t")
        if line[:1].isspace() or line[-1:].isspace():
            text_node.set(f"{{{XML}}}space", "preserve")
        text_node.text = line
        if template_end is not None:
            paragraph.append(copy.deepcopy(template_end))
        tx_body.append(paragraph)


def apply_replacement(root: ET.Element, slot_id: str, text: str, font_face: str | None = None) -> None:
    text_match = re.fullmatch(r"s\d{3}_sh(\d+)", slot_id)
    table_match = re.fullmatch(r"s\d{3}_tbl(\d+)_r(\d+)_c(\d+)", slot_id)
    if text_match:
        shape_id = text_match.group(1)
        for shape in root.findall(f".//{{{P}}}sp"):
            c_nv_pr = shape.find(f"{{{P}}}nvSpPr/{{{P}}}cNvPr")
            if c_nv_pr is not None and c_nv_pr.get("id") == shape_id:
                tx_body = shape.find(f"{{{P}}}txBody")
                if tx_body is None:
                    raise ValueError(f"text body missing for {slot_id}")
                replace_text_body(tx_body, text, font_face)
                return
    elif table_match:
        shape_id, row_index, col_index = map(int, table_match.groups())
        for frame in root.findall(f".//{{{P}}}graphicFrame"):
            c_nv_pr = frame.find(f"{{{P}}}nvGraphicFramePr/{{{P}}}cNvPr")
            if c_nv_pr is None or int(c_nv_pr.get("id", -1)) != shape_id:
                continue
            rows = frame.findall(f".//{{{A}}}tbl/{{{A}}}tr")
            if row_index >= len(rows):
                break
            cells = rows[row_index].findall(f"{{{A}}}tc")
            if col_index >= len(cells):
                break
            tx_body = cells[col_index].find(f"{{{A}}}txBody")
            if tx_body is None:
                raise ValueError(f"table text body missing for {slot_id}")
            replace_text_body(tx_body, text, font_face)
            return
    raise ValueError(f"slot not found in slide XML: {slot_id}")


def patch_cjk_runs(root: ET.Element, font_face: str) -> int:
    updated = 0
    for run in root.xpath(".//a:r | .//a:fld", namespaces={"a": A}):
        text = "".join(run.xpath("./a:t/text()", namespaces={"a": A}))
        if not any(ord(char) >= 0x2E80 for char in text):
            continue
        r_pr = run.find(f"{{{A}}}rPr")
        if r_pr is None:
            r_pr = ET.Element(f"{{{A}}}rPr")
            run.insert(0, r_pr)
        set_rpr_typefaces(r_pr, font_face)
        updated += 1
    return updated


def patch_theme_fonts(files: dict[str, bytes], font_face: str) -> None:
    for name in list(files):
        if not name.startswith("ppt/theme/theme") or not name.endswith(".xml"):
            continue
        root = ET.fromstring(files[name])
        for node in root.xpath(".//a:fontScheme/a:majorFont/a:ea | .//a:fontScheme/a:majorFont/a:cs | .//a:fontScheme/a:minorFont/a:ea | .//a:fontScheme/a:minorFont/a:cs", namespaces={"a": A}):
            node.set("typeface", font_face)
        files[name] = ET.tostring(root, encoding="UTF-8", xml_declaration=True, standalone=True)


def apply_plan(
    source: Path,
    library: dict[str, object],
    plan: dict[str, object],
    output: Path,
    font_face: str | None = None,
    font_scope: str = "replaced-text",
) -> None:
    if plan.get("status") != "confirmed":
        raise ValueError("fill plan must be confirmed before apply")
    report = check_plan(library, plan)
    if report["error_count"]:
        raise ValueError(f"fill plan has {report['error_count']} blocking error(s)")
    if source.resolve() == output.resolve():
        raise ValueError("output must not overwrite the source PPTX")
    with zipfile.ZipFile(source) as archive:
        files = {info.filename: archive.read(info.filename) for info in archive.infolist()}
        infos = archive.infolist()
    presentation, slide_ids, parts = presentation_slide_parts(files)
    slide_list = presentation.find(f"{{{P}}}sldIdLst")
    assert slide_list is not None
    for child in list(slide_list):
        slide_list.remove(child)
    for entry in plan["slides"]:
        source_slide = int(entry["source_slide"])
        slide_list.append(copy.deepcopy(slide_ids[source_slide - 1]))
        part = parts[source_slide - 1]
        root = ET.fromstring(files[part])
        for replacement in entry.get("replacements", []):
            apply_replacement(root, str(replacement["slot_id"]), str(replacement.get("text", "")), font_face)
        if font_face and font_scope == "all-selected-text":
            patch_cjk_runs(root, font_face)
        files[part] = ET.tostring(root, encoding="UTF-8", xml_declaration=True, standalone=True)
    if font_face and font_scope == "theme-and-replaced":
        patch_theme_fonts(files, font_face)
    files["ppt/presentation.xml"] = ET.tostring(presentation, encoding="UTF-8", xml_declaration=True, standalone=True)
    output.parent.mkdir(parents=True, exist_ok=True)
    with tempfile.NamedTemporaryFile(dir=output.parent, suffix=".pptx", delete=False) as handle:
        temp_path = Path(handle.name)
    try:
        with zipfile.ZipFile(temp_path, "w") as archive:
            for info in infos:
                archive.writestr(info, files[info.filename])
        shutil.move(str(temp_path), output)
    finally:
        temp_path.unlink(missing_ok=True)


def validate_output(pptx: Path, plan: dict[str, object]) -> dict[str, object]:
    errors: list[str] = []
    try:
        prs = Presentation(str(pptx))
    except Exception as exc:  # python-pptx surfaces malformed-package details
        return {"status": "FAIL", "errors": [f"pptx_open_failed:{exc}"], "slide_count": 0}
    if len(prs.slides) != len(plan.get("slides", [])):
        errors.append(f"slide_count_mismatch:{len(prs.slides)}:{len(plan.get('slides', []))}")
    for index, entry in enumerate(plan.get("slides", [])):
        if index >= len(prs.slides):
            break
        visible = "\n".join(shape.text for shape in prs.slides[index].shapes if getattr(shape, "has_text_frame", False))
        visible += "\n" + "\n".join(
            cell.text
            for shape in prs.slides[index].shapes if getattr(shape, "has_table", False)
            for row in shape.table.rows for cell in row.cells
        )
        for replacement in entry.get("replacements", []):
            value = str(replacement.get("text", ""))
            if value and value not in visible:
                errors.append(f"replacement_not_readable:slide={index + 1}:slot={replacement.get('slot_id')}")
    return {"status": "FAIL" if errors else "PASS", "errors": errors, "slide_count": len(prs.slides)}


def parse_slide_selection(value: str | None, slide_count: int) -> list[int]:
    if not value:
        return list(range(1, slide_count + 1))
    selected = [int(item.strip()) for item in value.split(",") if item.strip()]
    if any(item < 1 or item > slide_count for item in selected):
        raise ValueError("slide selection is out of range")
    if len(set(selected)) != len(selected):
        raise ValueError("repeating a source slide is not supported in v1")
    return selected


def main() -> None:
    parser = argparse.ArgumentParser(description=__doc__)
    sub = parser.add_subparsers(dest="command", required=True)
    analyze = sub.add_parser("analyze"); analyze.add_argument("source"); analyze.add_argument("--out", required=True)
    scaffold = sub.add_parser("scaffold"); scaffold.add_argument("library"); scaffold.add_argument("--slides"); scaffold.add_argument("--out", required=True)
    check = sub.add_parser("check-plan"); check.add_argument("library"); check.add_argument("plan"); check.add_argument("--out", required=True)
    apply = sub.add_parser("apply"); apply.add_argument("source"); apply.add_argument("library"); apply.add_argument("plan"); apply.add_argument("--out", required=True); apply.add_argument("--font"); apply.add_argument("--font-scope", choices=("replaced-text", "all-selected-text", "theme-and-replaced"))
    validate = sub.add_parser("validate"); validate.add_argument("pptx"); validate.add_argument("plan"); validate.add_argument("--out", required=True)
    args = parser.parse_args()

    if args.command == "analyze":
        result = analyze_pptx(Path(args.source).resolve()); write_json(Path(args.out), result)
    elif args.command == "scaffold":
        library = json.loads(Path(args.library).read_text(encoding="utf-8"))
        selected = parse_slide_selection(args.slides, int(library["slide_count"]))
        result = {
            "schema": "native_template_fill_plan.v1", "status": "draft",
            "source_pptx": library["source_pptx"],
            "font_policy": None,
            "slides": [{"source_slide": number, "purpose": "待填写", "layout_rationale": None, "replacements": []} for number in selected],
        }
        write_json(Path(args.out), result)
    elif args.command == "check-plan":
        library = json.loads(Path(args.library).read_text(encoding="utf-8")); plan = json.loads(Path(args.plan).read_text(encoding="utf-8"))
        result = check_plan(library, plan); write_json(Path(args.out), result)
        if result["error_count"]:
            raise SystemExit(2)
    elif args.command == "apply":
        source = Path(args.source).resolve(); library = json.loads(Path(args.library).read_text(encoding="utf-8")); plan = json.loads(Path(args.plan).read_text(encoding="utf-8"))
        font_policy = plan.get("font_policy") if isinstance(plan.get("font_policy"), dict) else {}
        font_face = args.font or font_policy.get("font_face")
        font_scope = args.font_scope or font_policy.get("scope") or "replaced-text"
        apply_plan(source, library, plan, Path(args.out).resolve(), font_face, font_scope)
        result = {"status": "PASS", "output": str(Path(args.out).resolve()), "font_face": font_face, "font_scope": font_scope}
    else:
        plan = json.loads(Path(args.plan).read_text(encoding="utf-8")); result = validate_output(Path(args.pptx).resolve(), plan); write_json(Path(args.out), result)
        if result["status"] != "PASS":
            raise SystemExit(2)
    print(json.dumps(result, ensure_ascii=False))


if __name__ == "__main__":
    main()
